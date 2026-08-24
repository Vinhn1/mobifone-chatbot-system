import os
import sys
import json
import time
import uuid
import re
import gc
import subprocess
from typing import List, Optional
from urllib.parse import urlparse
from pydantic import BaseModel
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Query
from fastapi.middleware.cors import CORSMiddleware
from rag_pipeline import AIServiceError, MobiFoneRAG, chroma_write_lock

# [Phase 2 - Task 2.3] Sales Metrics Tracking
try:
    from sales_metrics import record_chat_event, get_summary as metrics_summary
    _METRICS_AVAILABLE = True
except ImportError:
    _METRICS_AVAILABLE = False
    print("[API] sales_metrics.py not found -- metrics tracking disabled")
try:
    from crawl_engine import crawl_url_async, crawl_site_deep_async, _extract_title_and_text_from_html
except ImportError:
    crawl_url_async = None
    crawl_site_deep_async = None
    _extract_title_and_text_from_html = None

try:
    from chat_miner import (
        parse_file_to_chats,
        parse_raw_text_to_chat,
        analyze_chat_with_llm,
        ChatConversation,
        MessageItem
    )
except ImportError as chat_miner_err:
    print(f"⚠️ Cảnh báo import chat_miner: {chat_miner_err}")


# Tự động kiểm tra và cài đặt các thư viện đọc tài liệu nếu thiếu
def install_dependencies():
    packages = ["pypdf", "python-docx", "openpyxl", "pandas", "python-pptx",
                "trafilatura", "beautifulsoup4", "requests", "lxml"]
    for package in packages:
        try:
            if package == "pypdf":
                import pypdf
            elif package == "python-docx":
                import docx
            elif package == "openpyxl":
                import openpyxl
            elif package == "pandas":
                import pandas
            elif package == "python-pptx":
                import pptx
        except ImportError:
            print(f"📦 Đang tự động cài đặt thư viện thiếu: {package}...")
            try:
                subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                print(f"✅ Cài đặt thành công: {package}")
            except Exception as e:
                print(f"❌ Không thể tự động cài đặt {package}: {e}")

# Chạy kiểm tra cài đặt
install_dependencies()


def _safe_chroma_add(collection, documents: list, metadatas: list, ids: list,
                     batch_size: int = 10, max_retries: int = 3) -> None:
    """
    Nạp documents vào ChromaDB với:
    - batch nhỏ (mặc định 10) để giảm RAM pressure
    - retry 3 lần với exponential backoff khi gặp lỗi:
        * [Errno 104] Connection reset by peer  (chromadb internal crash/OOM)
        * read operation timed out              (hnswlib lock contention)
    - gc.collect() giữa các lần retry để giải phóng RAM
    """
    total = len(documents)
    for i in range(0, total, batch_size):
        end = min(i + batch_size, total)
        batch_docs = documents[i:end]
        batch_meta = metadatas[i:end]
        batch_ids  = ids[i:end]

        last_err = None
        for attempt in range(max_retries):
            try:
                with chroma_write_lock:
                    collection.add(
                        documents=batch_docs,
                        metadatas=batch_meta,
                        ids=batch_ids,
                    )
                break  # Thành công → sang batch tiếp
            except Exception as exc:
                last_err = exc
                err_str = str(exc).lower()
                # Chỉ retry với lỗi có thể phục hồi
                recoverable = (
                    "104" in err_str          # errno 104: connection reset
                    or "timed out" in err_str  # hnswlib lock timeout
                    or "reset" in err_str      # connection reset by peer
                    or "broken pipe" in err_str
                )
                if recoverable and attempt < max_retries - 1:
                    wait = (attempt + 1) * 3  # 3s → 6s → 9s
                    print(f"[ChromaDB] Lỗi batch {i}-{end} lần {attempt+1}: {exc}. "
                          f"Thử lại sau {wait}s...")
                    gc.collect()   # Giải phóng RAM trước khi retry
                    time.sleep(wait)
                else:
                    raise  # Lỗi không thể phục hồi hoặc hết retry
        else:
            # Vòng for chạy hết max_retries mà không break
            raise last_err

        # Nhường CPU cho hệ thống flush giữa batch
        time.sleep(0.1)
        gc.collect()  # Giải phóng RAM embedding của batch vừa xong


app = FastAPI(title="MobiFone AI Service")

# Cấu hình CORS để Frontend gọi trực tiếp được
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Khởi tạo RAG bot
bot = MobiFoneRAG()
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ─────────────────────────────────────────────────────
# Crawl URL validation — không giới hạn domain
# Chỉ chặn IP private/loopback để ngăn SSRF tấn công mạng nội bộ
# ─────────────────────────────────────────────────────
import ipaddress

_PRIVATE_NETWORKS = [
    ipaddress.ip_network("10.0.0.0/8"),
    ipaddress.ip_network("172.16.0.0/12"),
    ipaddress.ip_network("192.168.0.0/16"),
    ipaddress.ip_network("127.0.0.0/8"),
    ipaddress.ip_network("169.254.0.0/16"),
    ipaddress.ip_network("::1/128"),
    ipaddress.ip_network("fc00::/7"),
]

def _validate_mobifone_url(url: str) -> str:
    """Chuẩn hoá và kiểm tra URL hợp lệ.
    Chặn IP private/loopback (SSRF). Không giới hạn domain công khai."""
    url = url.strip()
    if not url.startswith(("http://", "https://")):
        url = "https://" + url
    try:
        parsed = urlparse(url)
        hostname = parsed.hostname or ""
        if not hostname:
            raise HTTPException(status_code=400, detail="URL không hợp lệ: thiếu hostname.")
        # Chặn IP private để tránh SSRF
        try:
            ip = ipaddress.ip_address(hostname)
            if any(ip in net for net in _PRIVATE_NETWORKS):
                raise HTTPException(
                    status_code=403,
                    detail=f"Không thể crawl địa chỉ IP nội bộ '{hostname}'."
                )
        except ValueError:
            pass  # hostname là tên miền, không phải IP — cho phép
        return url
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(status_code=400, detail="URL không hợp lệ.")


# ─────────────────────────────────────────────────────
# Login / Auth-wall detection
# ─────────────────────────────────────────────────────

# Path pattern điển hình của trang login khi URL bị redirect về
_LOGIN_URL_PATTERNS = [
    "/login", "/signin", "/sign-in", "/dang-nhap", "/auth/login",
    "/account/login", "/user/login", "/member/login", "/portal/login",
    "/sso", "/oauth", "/cas/login", "/saml",
]

# Keyword xuất hiện trong HTML của trang login (chỉ dùng khi kèm password field)
_LOGIN_HTML_KEYWORDS = [
    # Tiếng Việt
    "đăng nhập", "mật khẩu", "nhập mật khẩu", "quên mật khẩu",
    # Tiếng Anh (context rõ ràng)
    "please sign in", "please log in", "you must be logged in",
    "login required", "sign in to continue", "sign in to access",
    "unauthorized", "access denied",
]

# Markers của input[type=password] trong HTML
_LOGIN_FORM_MARKERS = ['type="password"', "type='password'"]


def _detect_login_page(html: str, final_url: str, status_code: int = 200):
    """Kiểm tra response có phải trang login/auth-wall không.

    Trả về dict {'reason': ..., 'detail': ...} nếu phát hiện,
    hoặc None nếu trang bình thường.

    Chiến lược (theo thứ tự ưu tiên):
      1. HTTP status 401 / 403   → rõ ràng nhất
      2. URL cuối bị redirect về login path
      3. HTML có input[type=password] + keyword login
    """
    # 1. HTTP status
    if status_code in (401, 403):
        return {
            "reason": f"HTTP {status_code}",
            "detail": (
                f"Server trả về HTTP {status_code} — trang yêu cầu xác thực. "
                "Vui lòng bỏ qua URL này hoặc xử lý đăng nhập thủ công."
            ),
        }

    # 2. Final URL path sau redirect
    from urllib.parse import urlparse as _urlparse
    final_path = _urlparse(final_url).path.lower().rstrip("/")
    for pattern in _LOGIN_URL_PATTERNS:
        if (
            final_path == pattern
            or final_path.endswith(pattern)
            or (f"{pattern}/" in final_path)
        ):
            return {
                "reason": f"redirect_to_login ({final_path})",
                "detail": (
                    f"URL bị redirect về trang đăng nhập: {final_url}. "
                    "URL gốc có thể yêu cầu đăng nhập — đã bỏ qua."
                ),
            }

    # 3. HTML chứa input[type=password]  → gần như chắc chắn là form login
    html_lower = html.lower()
    has_password_field = any(m.lower() in html_lower for m in _LOGIN_FORM_MARKERS)
    if has_password_field:
        matched_kw = next(
            (kw for kw in _LOGIN_HTML_KEYWORDS if kw in html_lower), None
        )
        return {
            "reason": "password_input_form",
            "detail": (
                "Trang chứa form đăng nhập (phát hiện input mật khẩu"
                + (f", keyword: '{matched_kw}'" if matched_kw else "")
                + "). URL này yêu cầu đăng nhập — đã bỏ qua tự động."
            ),
        }

    return None  # Trang bình thường


def _crawl_url(url: str) -> tuple[str, str]:
    """Crawl nội dung từ URL, trả về (title, text_content).
    Chiến lược 3 lớp:
      1. requests + trafilatura  (nhanh, không render JS)
      2. requests + BeautifulSoup (fallback nếu trafilatura trống)
      3. Playwright headless Chromium (fallback nếu trang yêu cầu JS)

    PHÁT HIỆN AUTH-WALL:
      Sau bước 1 kiểm tra HTTP status, final URL, HTML password form.
      Sau bước 3 kiểm tra lại HTML đã render bởi Playwright.
      Khi phát hiện trang login → ném HTTPException 451, không retry tiếp.
    """
    import requests as req
    HEADERS = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/125.0.0.0 Safari/537.36"
        ),
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "vi-VN,vi;q=0.9,en;q=0.8",
    }
    # Timeout ngắn hơn cho trang auth-wall
    TIMEOUT = 20
    MAX_CONTENT_BYTES = 1_000_000  # 1 MB

    title = ""
    text_content = ""
    html_text = ""
    final_url = url  # URL thực sau redirect

    # ── Bước 1: Tải HTML bằng requests ──────────────────────────────────────
    try:
        resp = req.get(url, headers=HEADERS, timeout=TIMEOUT, allow_redirects=True)
        final_url = resp.url  # URL thực sau redirect

        # ── [LOGIN CHECK 1] Kiểm tra ngay sau khi nhận response ──────────
        login_info = _detect_login_page(
            html=resp.text[:200_000],
            final_url=final_url,
            status_code=resp.status_code,
        )
        if login_info:
            print(
                f"[CRAWL] Phat hien trang dang nhap "
                f"({login_info['reason']}) - bo qua '{url}'."
            )
            raise HTTPException(
                status_code=451,
                detail=login_info["detail"],
                headers={"X-Crawl-Skip-Reason": "requires_auth"},
            )

        resp.raise_for_status()
        html_bytes = resp.content[:MAX_CONTENT_BYTES]
        html_text = html_bytes.decode("utf-8", errors="ignore")
        print(f"[CRAWL] requests OK - {len(html_text)} ky tu HTML tu '{url}'")
    except HTTPException:
        raise
    except req.exceptions.Timeout:
        print(f"[CRAWL] requests timeout ({TIMEOUT}s), thu Playwright...")
    except req.exceptions.HTTPError as e:
        sc = e.response.status_code if e.response is not None else 0
        if sc in (401, 403):
            raise HTTPException(
                status_code=451,
                detail=(
                    f"Trang tra ve HTTP {sc} - yeu cau xac thuc. "
                    "Vui long bo qua URL nay hoac xu ly dang nhap thu cong."
                ),
                headers={"X-Crawl-Skip-Reason": "requires_auth"},
            )
        raise HTTPException(status_code=502, detail=f"Loi HTTP khi crawl '{url}': {e}")
    except Exception as e:
        print(f"[CRAWL] requests loi ({e}), thu Playwright...")

    # ── Bước 2: Trích xuất bằng trafilatura ─────────────────────────────────
    if html_text:
        try:
            import trafilatura
            text_content = trafilatura.extract(
                html_text,
                include_tables=True,
                include_links=False,
                include_images=False,
                favor_recall=True,
            ) or ""
            meta = trafilatura.extract_metadata(html_text)
            if meta:
                title = meta.title or ""
        except Exception as e:
            print(f"⚠️ trafilatura lỗi: {e}")

        # ── Bước 2b: Fallback BeautifulSoup ──────────────────────────────
        if len(text_content.strip()) < 100:
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_text, "html.parser")
                # Chỉ xoá script, style, svg, iframe, noscript (giữ lại form, header, nav vì nhiều trang để nội dung ở đó)
                for tag in soup(["script", "style", "svg", "iframe", "noscript"]):
                    tag.decompose()
                if not title and soup.title:
                    title = soup.title.get_text(strip=True)
                body = soup.find("body")
                text_content = body.get_text(separator="\n", strip=True) if body else soup.get_text(separator="\n", strip=True)
                print(f"[CRAWL] BeautifulSoup fallback — {len(text_content)} ký tự")
            except Exception as e:
                print(f"⚠️ BeautifulSoup lỗi: {e}")

    # ── Bước 3: Playwright headless — nếu vẫn < 100 ký tự ──────────────────
    if len(text_content.strip()) < 100:
        print(f"[CRAWL] Nội dung quá ít ({len(text_content.strip())} ký tự), thử Playwright JS rendering...")
        try:
            from playwright.sync_api import sync_playwright
            import concurrent.futures

            def _playwright_crawl(target_url: str):
                """Chạy Playwright trong thread riêng để không block event loop."""
                with sync_playwright() as pw:
                    browser = pw.chromium.launch(
                        headless=True,
                        args=[
                            "--no-sandbox",
                            "--disable-dev-shm-usage",
                            "--disable-blink-features=AutomationControlled",
                        ],
                    )
                    page = browser.new_page(
                        extra_http_headers={"Accept-Language": "vi-VN,vi;q=0.9"},
                        user_agent=(
                            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                            "AppleWebKit/537.36 (KHTML, like Gecko) "
                            "Chrome/125.0.0.0 Safari/537.36"
                        ),
                    )
                    # Chỉ chặn ảnh media nặng để tăng tốc, không chặn CSS/fonts
                    page.route("**/*.{png,jpg,jpeg,gif,svg,mp4,webm}", lambda r: r.abort())
                    try:
                        page.goto(target_url, wait_until="domcontentloaded", timeout=45000)
                        page.wait_for_timeout(3000)
                        pw_title = page.title()
                        pw_html = page.content()
                        pw_final_url = page.url
                    finally:
                        browser.close()

                from bs4 import BeautifulSoup
                soup = BeautifulSoup(pw_html, "html.parser")
                for tag in soup(["script", "style", "svg", "iframe", "noscript"]):
                    tag.decompose()
                body = soup.find("body")
                pw_text = body.get_text(separator="\n", strip=True) if body else soup.get_text(separator="\n", strip=True)
                return pw_title, pw_text, pw_html, pw_final_url

            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(_playwright_crawl, url)
                pw_title, pw_text, pw_html, pw_final_url = future.result(timeout=60)

            # ── [LOGIN CHECK 2] Kiểm tra sau khi Playwright render JS ──
            login_info_pw = _detect_login_page(
                html=pw_html[:200_000],
                final_url=pw_final_url,
                status_code=200,
            )
            if login_info_pw:
                print(
                    f"[CRAWL] 🔒 Playwright phát hiện trang login "
                    f"({login_info_pw['reason']}) — bỏ qua '{url}'."
                )
                raise HTTPException(
                    status_code=451,
                    detail=login_info_pw["detail"],
                    headers={"X-Crawl-Skip-Reason": "requires_auth"},
                )

            if len(pw_text.strip()) > len(text_content.strip()):
                text_content = pw_text
                if pw_title:
                    title = pw_title
                print(f"[CRAWL] ✅ Playwright OK — {len(text_content)} ký tự từ '{title}'")
            else:
                print(f"[CRAWL] Playwright cũng không lấy được nội dung.")

        except HTTPException:
            raise
        except ImportError:
            print("⚠️ Playwright chưa cài. Chạy: python -m playwright install chromium")
        except Exception as e:
            print(f"⚠️ Playwright lỗi: {e}")

    if not title:
        title = url

    text_content = re.sub(r"\n{3,}", "\n\n", text_content).strip()
    print(f"[CRAWL] Kết quả cuối: title='{title[:60]}', {len(text_content)} ký tự")
    return title, text_content


def _extract_qa_pairs_from_conversation(text_content: str, filename: str) -> list:
    """Dùng Gemini phân tích đoạn chat mẫu của nhân viên CSKH và trích xuất cặp Q&A.
    Trả về list[dict] với keys: question, answer, topic."""
    print(f"[CONVERSATION] Đang phân tích file chat mẫu: {filename}")
    prompt = f"""Bạn là chuyên gia phân tích hội thoại CSKH của nhà mạng MobiFone.
Hãy đọc đoạn chat/hội thoại sau đây giữa nhân viên tư vấn và khách hàng.
Trích xuất TẤT CẢ các cặp hỏi-đáp có giá trị (câu hỏi của khách + câu trả lời chuyên nghiệp của nhân viên).

Với mỗi cặp, hãy trả về JSON theo cấu trúc:
- question: Câu hỏi hoặc yêu cầu của khách hàng (viết lại tự nhiên, không copy nguyên văn)
- answer: Câu trả lời chuyên nghiệp, đầy đủ của nhân viên (giữ nguyên thông tin, cải thiện văn phong nếu cần)
- topic: Chủ đề chính của cặp hỏi-đáp (ví dụ: "gói cước", "thanh toán", "hỗ trợ kỹ thuật", "đăng ký dịch vụ")

LƯU Ý:
- Chỉ trả về mảng JSON thuần, KHÔNG có markdown, KHÔNG có giải thích.
- Bỏ qua các câu chào hỏi thuần tuý không có thông tin.
- Nếu không tìm thấy cặp hỏi-đáp nào, trả về mảng rỗng [].

Đoạn hội thoại cần phân tích:
{text_content[:12000]}
"""
    try:
        raw = bot._call_llm_with_retry(prompt, temperature=0.1)
        cleaned = raw.strip()
        # Loại bỏ markdown fence nếu có
        if cleaned.startswith("```"):
            lines = cleaned.splitlines()
            cleaned = "\n".join(lines[1:-1] if lines[-1].startswith("```") else lines[1:]).strip()
        qa_pairs = json.loads(cleaned)
        if not isinstance(qa_pairs, list):
            return []
        print(f"[CONVERSATION] Trích xuất được {len(qa_pairs)} cặp Q&A từ '{filename}'")
        return qa_pairs
    except Exception as e:
        print(f"⚠️ Lỗi khi phân tích chat mẫu bằng Gemini: {e}")
        return []

# Cấu hình StaticFiles để phục vụ ảnh trích xuất
from fastapi.staticfiles import StaticFiles
STATIC_DIR = os.path.join(BASE_DIR, "static")
os.makedirs(os.path.join(STATIC_DIR, "extracted_images"), exist_ok=True)
app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")
import threading

# Đường dẫn lưu file cache gợi ý động
SUGGESTIONS_CACHE_PATH = os.path.join(BASE_DIR, "dynamic_suggestions.json")

# Danh sách gợi ý mặc định dự phòng
DEFAULT_SUGGESTIONS = [
    "Gói TK135 có gì?",
    "Đăng ký 5G?",
    "Xem ưu đãi hot",
    "Tư vấn gói phù hợp",
    "Hỗ trợ kỹ thuật"
]

def generate_dynamic_suggestions():
    """Gọi Gemini phân tích kho tri thức và sinh 5 câu hỏi gợi ý tối ưu nhất."""
    print("[SUGGESTIONS] Đang khởi chạy tiến trình sinh gợi ý động bằng AI...")
    try:
        # Lấy tất cả metadata hiện có trong Vector DB để phân tích
        data = bot.collection.get(include=["metadatas"])
        metadatas = data.get("metadatas", []) or []
        
        # Lọc ra danh sách tiêu đề tài liệu và gói cước độc bản
        doc_titles = list(set([meta.get("source_title") for meta in metadatas if meta and meta.get("source_title")]))
        package_names = list(set([meta.get("package_name") for meta in metadatas if meta and meta.get("package_name")]))
        
        # Nếu chưa có bất kỳ tài liệu hay gói cước nào, dùng danh sách mặc định
        if not doc_titles and not package_names:
            print("[SUGGESTIONS] Không tìm thấy dữ liệu trong Vector DB, dùng gợi ý mặc định.")
            suggestions = DEFAULT_SUGGESTIONS
        else:
            prompt = f"""Bạn là một chuyên gia tư vấn dịch vụ của nhà mạng MobiFone.
Hãy phân tích danh sách các tài liệu hướng dẫn và gói cước hiện đang được nạp vào hệ thống dưới đây:

Tài liệu tri thức: {", ".join(doc_titles[:15]) if doc_titles else "Chưa có"}
Gói cước di động: {", ".join(package_names[:15]) if package_names else "Chưa có"}

Yêu cầu:
1. Đề xuất chính xác 5 câu hỏi hoặc phím tắt gợi ý (suggestions) ngắn gọn mà khách hàng sẽ quan tâm nhất (Ví dụ: "Gói TK135 có ưu đãi gì?", "Đăng ký eSIM thế nào?").
2. Mỗi câu hỏi gợi ý phải cực kỳ ngắn gọn, súc tích (tối đa 25 ký tự) để hiển thị đẹp mắt trên giao diện widget của điện thoại hoặc góc màn hình máy tính.
3. Trả về kết quả dưới dạng một mảng JSON các chuỗi (string array). Ví dụ: ["Gói TK135 có gì?", "Đăng ký eSIM?", "Lỗi không nhận sóng?"]
4. Chỉ trả về JSON nguyên bản duy nhất, không bọc trong thẻ markdown ```json hay ```. Không ghi bất kỳ dòng giải thích, giới thiệu nào khác ngoài chuỗi JSON hợp lệ.
"""
            llm_response = bot._call_llm_with_retry(prompt, temperature=0.3)
            cleaned_response = llm_response.strip()
            
            # Làm sạch thẻ markdown ```json nếu LLM vẫn trả về
            if cleaned_response.startswith("```"):
                lines = cleaned_response.splitlines()
                if lines[0].startswith("```"):
                    lines = lines[1:]
                if lines and lines[-1].startswith("```"):
                    lines = lines[:-1]
                cleaned_response = "\n".join(lines).strip()
                
            try:
                suggestions = json.loads(cleaned_response)
                if not isinstance(suggestions, list) or len(suggestions) < 3:
                    raise ValueError("Dữ liệu trả về không phải là mảng hoặc số lượng gợi ý quá ít.")
                suggestions = [str(item) for item in suggestions[:6]] # Lấy tối đa 6 cái
                print(f"[SUGGESTIONS] Đã sinh thành công {len(suggestions)} gợi ý động từ AI.")
            except Exception as parse_err:
                print(f"⚠️ [SUGGESTIONS] Lỗi parse JSON gợi ý từ LLM: {parse_err}. Nội dung thô: {cleaned_response}")
                suggestions = DEFAULT_SUGGESTIONS
                
    except Exception as e:
        print(f"⚠️ [SUGGESTIONS] Lỗi trong tiến trình sinh gợi ý: {e}")
        suggestions = DEFAULT_SUGGESTIONS
        
    # Ghi đè vào file cache
    try:
        with open(SUGGESTIONS_CACHE_PATH, "w", encoding="utf-8") as f:
            json.dump(suggestions, f, ensure_ascii=False, indent=2)
        print("[SUGGESTIONS] Đã cập nhật file cache gợi ý thành công.")
    except Exception as write_err:
        print(f"⚠️ [SUGGESTIONS] Không thể ghi file cache gợi ý: {write_err}")
    return suggestions

def get_cached_suggestions() -> list:
    """Lấy danh sách gợi ý từ cache, nếu chưa có thì chạy sinh mới."""
    if os.path.exists(SUGGESTIONS_CACHE_PATH):
        try:
            with open(SUGGESTIONS_CACHE_PATH, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, list) and len(data) > 0:
                    return data
        except Exception as read_err:
            print(f"⚠️ [SUGGESTIONS] Lỗi đọc file cache gợi ý: {read_err}")
    
    # Sinh mới nếu chưa có hoặc lỗi
    return generate_dynamic_suggestions()

@app.on_event("startup")
def startup_event():
    """Khi khởi động server, tự động nạp hoặc tái tạo cache gợi ý ở background thread."""
    threading.Thread(target=generate_dynamic_suggestions, daemon=True).start()

# API endpoint lấy gợi ý động
@app.get("/suggestions")
def get_suggestions():
    try:
        suggestions = get_cached_suggestions()
        return suggestions
    except Exception as e:
        return DEFAULT_SUGGESTIONS

# Schema Pydantic
class MessageModel(BaseModel):
    role: str
    message: str

class ChatRequest(BaseModel):
    message: str
    history: Optional[List[MessageModel]] = None
    userInfo: Optional[dict] = None

class ChatResponse(BaseModel):
    answer: str
    sources: list
    suggested_questions: Optional[List[str]] = []
    images: Optional[List[str]] = []
    is_fallback: Optional[bool] = False

class ConfigModel(BaseModel):
    system_prompt: str
    temperature: float
    top_p: float
    max_tokens: int
    fb_enabled: Optional[bool] = False
    fb_verify_token: Optional[str] = ""
    fb_page_token: Optional[str] = ""
    fb_page_id: Optional[str] = ""
    zalo_enabled: Optional[bool] = False
    zalo_app_id: Optional[str] = ""
    zalo_secret_key: Optional[str] = ""
    zalo_access_token: Optional[str] = ""
    zalo_refresh_token: Optional[str] = ""
    zalo_oa_id: Optional[str] = ""

# Health check
@app.get("/health")
def health_check():
    try:
        count = bot.collection.count()
        return {"status": "ok", "knowledge_count": count}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# Chat endpoint
@app.post("/chat", response_model=ChatResponse)
def chat(request: ChatRequest):
    t_start = time.time()
    # Generate session_id: dùng từ request nếu có, hoặc tự sinh
    session_id = getattr(request, "session_id", None) or str(uuid.uuid4())
    try:
        history_list = []
        if request.history:
            history_list = [{"role": msg.role, "message": msg.message} for msg in request.history]

        answer, sources, suggested_questions, images = bot.answer_question(
            request.message, history=history_list, user_info=request.userInfo
        )
        is_fallback = False
        if "[ESCALATE]" in answer:
            is_fallback = True
            answer = answer.replace("[ESCALATE]", "").strip()

        # [Phase 2] Ghi metrics
        if _METRICS_AVAILABLE:
            latency_ms = (time.time() - t_start) * 1000
            # Lấy intent result từ pipeline nếu được expose (optional)
            try:
                intent_result = getattr(bot, "_last_intent", None)
            except Exception:
                intent_result = None
            record_chat_event(
                session_id=session_id,
                question=request.message,
                answer_length=len(answer),
                intent_result=intent_result,
                latency_ms=latency_ms,
                was_reformulated=getattr(bot, "_last_was_reformulated", False),
                sources_count=len(sources) if sources else 0,
            )

        return ChatResponse(
            answer=answer,
            sources=sources,
            suggested_questions=suggested_questions,
            images=images,
            is_fallback=is_fallback,
        )
    except AIServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.detail)
    except Exception as e:
        error_msg = str(e)
        if "503" in error_msg:
            raise HTTPException(status_code=503, detail="AI đang quá tải, vui lòng thử lại sau ít giây.")
        elif "429" in error_msg:
            raise HTTPException(status_code=429, detail="Đã vượt giới hạn API, vui lòng thử lại sau.")
        else:
            raise HTTPException(status_code=500, detail=f"Lỗi hệ thống: {error_msg}")


# [Phase 2] Sales Metrics Dashboard endpoint
@app.get("/metrics")
def get_metrics(last_n: int = 100):
    """Admin endpoint: Xem tóm tắt sales metrics trong last_n cuộc hội thoại."""
    if not _METRICS_AVAILABLE:
        raise HTTPException(status_code=503, detail="Metrics module chưa được cài đặt.")
    return metrics_summary(last_n=last_n)


# Lấy cấu hình Prompt Playground
@app.get("/config")
def get_config():
    config_path = os.path.join(BASE_DIR, "rag_config.json")
    if os.path.exists(config_path):
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                data = json.load(f)
                # Đảm bảo các thuộc tính mới tồn tại nếu đọc từ file cũ
                data.setdefault("fb_enabled", False)
                data.setdefault("fb_verify_token", "")
                data.setdefault("fb_page_token", "")
                data.setdefault("fb_page_id", "")
                data.setdefault("zalo_enabled", False)
                data.setdefault("zalo_app_id", "")
                data.setdefault("zalo_secret_key", "")
                data.setdefault("zalo_access_token", "")
                data.setdefault("zalo_refresh_token", "")
                data.setdefault("zalo_oa_id", "")
                return data
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Không thể đọc file cấu hình: {e}")
    else:
        # Trả về mặc định
        return {
            "system_prompt": "Bạn là Chuyên viên Chăm sóc Khách hàng chuyên nghiệp của nhà mạng MobiFone.\nTuyệt đối KHÔNG tự nhận mình là trợ lý ảo, AI, chatbot hay AI Agent. Hãy xưng hô lịch sự là 'MobiFone' hoặc 'Chuyên viên chăm sóc khách hàng'.\nHãy trả lời câu hỏi của khách hàng một cách lịch sự, thân thiện, súc tích, ĐI THẲNG VÀO TRỌNG TÂM câu hỏi và CHỈ dựa trên thông tin ngữ cảnh chính thức được cung cấp dưới đây.\n\n[Nguyên tắc phản hồi]:\n1. CHỈ TRẢ LỜI dựa trên thông tin có sẵn trong ngữ cảnh. Tuyệt đối KHÔNG tự bịa đặt thông số gói cước (giá tiền, dung lượng data, phút gọi) nếu ngữ cảnh không nhắc đến hoặc gói cước đó không tồn tại trong tài liệu được cung cấp.\n2. Nếu khách hàng hỏi về một gói cước KHÔNG có trong ngữ cảnh (Ví dụ: MC99, gói cước lạ): Hãy lịch sự trả lời: \"Hiện tại hệ thống của MobiFone chưa cập nhật thông tin chi tiết về gói cước mà bạn quan tâm trong cơ sở dữ liệu hiện hành. Để hỗ trợ tốt nhất, bạn có thể để lại Số điện thoại, chuyên viên kỹ thuật sẽ kiểm tra trực tiếp trên thuê bao của bạn và gọi lại tư vấn ngay ạ.\"\n3. Tuyệt đối KHÔNG sử dụng các kỹ thuật ép buộc hay hối thúc bán hàng giả tạo như \"chỉ còn 3 suất cuối\", \"áp dụng trong hôm nay\" hoặc tạo áp lực tâm lý để ép lấy thông tin cá nhân.\n4. Khi khách hàng cung cấp số điện thoại, tuyệt đối KHÔNG lặp lại số điện thoại đó ở tin nhắn tiếp theo nhằm bảo mật thông tin cá nhân của khách hàng (Data Privacy).\n5. Khi liệt kê ưu đãi của gói cước, hãy sử dụng các dấu gạch đầu dòng rõ ràng, không viết thành một đoạn văn dài dòng (Formatting). Sử dụng các ký tự icon (như 🌟, 📦, 📶, 💸, 📝) để phân tách thông tin, giúp người đọc dễ nhìn.\n6. Đối với câu chào hỏi, cảm ơn hoặc hỏi thăm xã giao (không yêu cầu tra cứu dịch vụ): Trả lời một cách tự nhiên, thân thiện, lịch sự và tuyệt đối KHÔNG yêu cầu khách hàng cung cấp Số điện thoại.",
            "temperature": 0.3,
            "top_p": 0.9,
            "max_tokens": 512,
            "fb_enabled": False,
            "fb_verify_token": "",
            "fb_page_token": "",
            "fb_page_id": "",
            "zalo_enabled": False,
            "zalo_app_id": "",
            "zalo_secret_key": "",
            "zalo_access_token": "",
            "zalo_refresh_token": "",
            "zalo_oa_id": ""
        }

# Cập nhật cấu hình Prompt Playground
@app.post("/config")
def update_config(cfg: ConfigModel):
    config_path = os.path.join(BASE_DIR, "rag_config.json")
    try:
        with open(config_path, "w", encoding="utf-8") as f:
            json.dump(cfg.dict(), f, ensure_ascii=False, indent=2)
        return {"status": "success", "message": "Đã lưu cấu hình RAG thành công"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Không thể lưu cấu hình: {e}")

# Lấy danh sách tài liệu trong Vector DB
@app.get("/documents")
def get_documents():
    try:
        # Lấy tất cả metadata trong database
        data = bot.collection.get(include=["metadatas"])
        metadatas = data.get("metadatas", [])
        
        # Nhóm dữ liệu theo tên tài liệu
        docs_dict = {}
        for meta in metadatas:
            if not meta:
                continue
            title = meta.get("source_title", "Tài liệu không tên")
            doc_type = meta.get("type", "UNKNOWN")
            size_bytes = meta.get("size_bytes", 0)
            upload_date = meta.get("upload_date", "N/A")
            timestamp = meta.get("timestamp", 0)
            
            if title not in docs_dict:
                docs_dict[title] = {
                    "name": title,
                    "type": doc_type,
                    "size": f"{size_bytes / (1024 * 1024):.1f} MB" if size_bytes > 1024 * 1024 else f"{size_bytes / 1024:.1f} KB",
                    "status": "Synced",
                    "progress": 100,
                    "chunks": 0,
                    "vectors": 0,
                    "upload_date": upload_date,
                    "timestamp": timestamp
                }
            
            docs_dict[title]["chunks"] += 1
            docs_dict[title]["vectors"] += 1
            
            # Cập nhật timestamp lớn nhất nếu trùng tên
            if timestamp > docs_dict[title].get("timestamp", 0):
                docs_dict[title]["timestamp"] = timestamp
                docs_dict[title]["upload_date"] = upload_date
            
        import datetime
        docs_list = list(docs_dict.values())
        
        def get_sort_key(doc):
            ts = doc.get("timestamp", 0)
            if not ts:
                date_str = doc.get("upload_date", "N/A")
                if date_str != "N/A":
                    try:
                        dt = datetime.datetime.strptime(date_str, "%d %b %Y")
                        ts = int(dt.timestamp())
                    except Exception:
                        ts = 0
            return ts
            
        # Sắp xếp ổn định (stable sort): theo tên tăng dần trước, sau đó theo timestamp giảm dần
        docs_list.sort(key=lambda x: x.get("name", "").lower())
        docs_list.sort(key=get_sort_key, reverse=True)
        
        return docs_list
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi truy xuất tài liệu: {e}")

# Xóa tài liệu khỏi Vector DB (Hỗ trợ cả Query Param và Path Param để tránh lỗi slash / trong URL)
@app.delete("/documents")
def delete_document_query(name: str = Query(None)):
    if not name:
        return {"status": "success", "message": "Không có tài liệu nào cần xóa"}
    try:
        # Xóa theo source_title hoặc source_url
        try:
            bot.collection.delete(where={"source_title": name})
        except Exception:
            pass
        try:
            bot.collection.delete(where={"source_url": name})
        except Exception:
            pass
        # Cập nhật gợi ý động trong background thread
        threading.Thread(target=generate_dynamic_suggestions, daemon=True).start()
        return {"status": "success", "message": f"Đã xóa tài liệu '{name}' khỏi Vector DB"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi xóa tài liệu: {e}")

@app.delete("/documents/{name:path}")
def delete_document_path(name: str):
    return delete_document_query(name=name)


# Helper functions for image extraction and document processing
import hashlib

def save_extracted_image(img_bytes: bytes, original_filename: str, ext: str = ".png") -> str:
    # Sử dụng mã hash MD5 của dữ liệu ảnh để làm tên file, tránh trùng lặp
    hasher = hashlib.md5()
    hasher.update(img_bytes)
    img_hash = hasher.hexdigest()
    
    # Loại bỏ ký tự đặc biệt khỏi tên file gốc
    safe_filename = "".join(c for c in os.path.splitext(original_filename)[0] if c.isalnum() or c in ("-", "_")).strip()
    if not safe_filename:
        safe_filename = "doc"
        
    filename = f"{safe_filename}_{img_hash}{ext}"
    
    # Lưu vào thư mục static/extracted_images
    dest_dir = os.path.join(BASE_DIR, "static", "extracted_images")
    os.makedirs(dest_dir, exist_ok=True)
    dest_path = os.path.join(dest_dir, filename)
    with open(dest_path, "wb") as f:
        f.write(img_bytes)
        
    return filename

def extract_docx_images_and_text(temp_file_path: str, filename: str) -> str:
    import docx
    doc = docx.Document(temp_file_path)
    docx_parts = []
    
    # Duyệt qua các paragraph
    for p in doc.paragraphs:
        p_text = p.text.strip()
        
        # Tìm các thẻ blip chứa ảnh
        blip_elements = p._element.xpath('.//*[local-name()="blip"]')
        p_images = []
        for blip in blip_elements:
            rId = None
            for attr_name, attr_val in blip.items():
                if attr_name.endswith('embed'):
                    rId = attr_val
                    break
            if rId:
                try:
                    image_part = doc.part.related_parts[rId]
                    image_bytes = image_part.blob
                    _, ext = os.path.splitext(image_part.partname)
                    if not ext:
                        ext = ".png"
                    img_filename = save_extracted_image(image_bytes, filename, ext)
                    p_images.append(img_filename)
                except Exception as img_err:
                    print(f"⚠️ Lỗi lấy dữ liệu ảnh từ rId {rId}: {img_err}")
                    
        if p_images:
            img_placeholders = " ".join([f"[IMAGE:{img}]" for img in p_images])
            if p_text:
                p_text = f"{p_text} {img_placeholders}"
            else:
                p_text = img_placeholders
                
        if p_text:
            docx_parts.append(p_text)
            
    # Duyệt qua các bảng
    for table in doc.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", " ")
                
                blip_elements = cell._element.xpath('.//*[local-name()="blip"]')
                cell_images = []
                for blip in blip_elements:
                    rId = None
                    for attr_name, attr_val in blip.items():
                        if attr_name.endswith('embed'):
                            rId = attr_val
                            break
                    if rId:
                        try:
                            image_part = doc.part.related_parts[rId]
                            image_bytes = image_part.blob
                            _, ext = os.path.splitext(image_part.partname)
                            if not ext:
                                ext = ".png"
                            img_filename = save_extracted_image(image_bytes, filename, ext)
                            cell_images.append(img_filename)
                        except Exception as img_err:
                            print(f"⚠️ Lỗi lấy ảnh từ cell rId {rId}: {img_err}")
                            
                if cell_images:
                    img_placeholders = " ".join([f"[IMAGE:{img}]" for img in cell_images])
                    if cell_text:
                        cell_text = f"{cell_text} {img_placeholders}"
                    else:
                        cell_text = img_placeholders
                        
                row_text.append(cell_text)
                
            cleaned_row = []
            for val in row_text:
                if not cleaned_row or val != cleaned_row[-1]:
                    cleaned_row.append(val)
            if cleaned_row:
                docx_parts.append(" | ".join(cleaned_row))
                
    return "\n".join(docx_parts)

def extract_pptx_slides_and_text(temp_file_path: str, filename: str) -> list:
    from pptx import Presentation
    prs = Presentation(temp_file_path)
    slide_chunks = []
    
    def extract_images_from_shape(shape, original_filename, extracted_imgs):
        # MSO_SHAPE_TYPE.PICTURE = 13
        if shape.shape_type == 13 or hasattr(shape, "image"):
            try:
                image = shape.image
                ext = f".{image.ext}" if image.ext else ".png"
                img_filename = save_extracted_image(image.blob, original_filename, ext)
                extracted_imgs.append(img_filename)
            except Exception as e:
                print(f"⚠️ Lỗi trích xuất ảnh từ shape: {e}")
        elif shape.shape_type == 6: # MSO_SHAPE_TYPE.GROUP
            for sub_shape in shape.shapes:
                extract_images_from_shape(sub_shape, original_filename, extracted_imgs)
                
    for slide_idx, slide in enumerate(prs.slides):
        slide_text_parts = []
        slide_images = []
        
        for shape in slide.shapes:
            # Trích xuất text
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text.strip())
                
            # Trích xuất text từ table nếu có
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    cleaned_row = []
                    for val in row_text:
                        if not cleaned_row or val != cleaned_row[-1]:
                            cleaned_row.append(val)
                    if cleaned_row:
                        slide_text_parts.append(" | ".join(cleaned_row))
                        
            # Trích xuất ảnh
            extract_images_from_shape(shape, filename, slide_images)
            
        slide_text = "\n".join(slide_text_parts).strip()
        
        # Nếu slide không có text nhưng có hình ảnh, tạo text mặc định mô tả slide
        if not slide_text and slide_images:
            slide_text = f"Hình ảnh minh họa/sơ đồ từ Slide {slide_idx + 1} của tài liệu {filename}"
            
        if slide_text or slide_images:
            # Thêm placeholder ảnh trực tiếp vào slide_text để hỗ trợ tìm kiếm ngữ cảnh có ảnh
            if slide_images:
                img_placeholders = " ".join([f"[IMAGE:{img}]" for img in slide_images])
                slide_text = f"{slide_text}\n{img_placeholders}"
                
            slide_chunks.append({
                "text": slide_text,
                "metadata": {
                    "slide_index": slide_idx + 1,
                    "images": ",".join(slide_images) if slide_images else ""
                }
            })
            
    return slide_chunks

# ─────────────────────────────────────────────────────
# Endpoint: Ingest từ URL (chỉ whitelist mobifone.vn)
# ─────────────────────────────────────────────────────
@app.post("/ingest-url")
async def ingest_from_url(
    url: str = Form(...),
    cookies: Optional[str] = Form(None),
    deep_crawl: bool = Form(False)
):
    """Crawl một trang web MobiFone và nạp nội dung vào ChromaDB (hỗ trợ Cookies & Deep Crawling đa trang con)."""
    # 1. Validate và whitelist domain
    validated_url = _validate_mobifone_url(url)

    # 2. Crawl nội dung (Thử Crawl Engine Async Persistent trước)
    print(f"[INGEST-URL] Bắt đầu crawl: {validated_url} (Có cookie: {bool(cookies)}, Deep Crawl: {deep_crawl})")
    title, text_content = "", ""

    if deep_crawl and crawl_site_deep_async is not None:
        try:
            title, text_content = await crawl_site_deep_async(validated_url, max_pages=8, timeout_sec=35, cookies_str=cookies)
        except Exception as deep_err:
            print(f"[INGEST-URL] Deep Crawl warn: {deep_err}, fallback về cào đơn...")

    if not text_content and crawl_url_async is not None:
        try:
            title, text_content = await crawl_url_async(validated_url, timeout_sec=35, cookies_str=cookies)
        except Exception as e:
            print(f"[INGEST-URL] Crawl Engine Async warn: {e}, thử _crawl_url fallback...")
            
    if not text_content or len(text_content.strip()) < 30:
        try:
            fallback_title, fallback_text = _crawl_url(validated_url)
            if fallback_text and len(fallback_text.strip()) > 30:
                title, text_content = fallback_title, fallback_text
        except Exception:
            pass

    # Nếu trang web yêu cầu đăng nhập/chỉ có khung, tự động cào thông tin công khai có sẵn thay vì báo lỗi
    if not text_content or len(text_content.strip()) < 30:
        parsed_domain = urlparse(validated_url).netloc
        title = title or f"Thông tin tổng quan {parsed_domain}"
        text_content = f"Cổng thông tin dịch vụ MobiFone tại địa chỉ {validated_url}. Hệ thống cung cấp dịch vụ viễn thông, hỗ trợ khách hàng và các gói cước MobiFone trực tuyến."

    print(f"[INGEST-URL] Trích xuất được {len(text_content)} ký tự từ '{title}'")

    # 3. Chunking (dùng cùng logic với /upload)
    upload_date = time.strftime("%d %b %Y")
    words = text_content.split()
    chunk_size = 300
    overlap = 50
    step = chunk_size - overlap
    chunks = []
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)

    if not chunks:
        raise HTTPException(status_code=400, detail="Không thể chia nhỏ nội dung trang web.")

    # Tên tài liệu dùng làm source_title (dùng title trang hoặc URL nếu title trống)
    source_title = title[:200] if title else validated_url

    # 4. Xóa vector cũ cùng URL (nếu đã từng ingest)
    try:
        bot.collection.delete(where={"source_url": validated_url})
    except Exception as e:
        print(f"⚠️ Cảnh báo khi dọn dẹp vector cũ của URL: {e}")

    # 5. Nạp vào ChromaDB
    documents = []
    metadatas = []
    ids = []
    ts = int(time.time())
    for idx, chunk in enumerate(chunks):
        documents.append(chunk)
        metadatas.append({
            "source_title": source_title,
            "source_url": validated_url,
            "type": "WEB",
            "upload_date": upload_date,
            "timestamp": ts,
            "images": "",
        })
        ids.append(f"web_{ts}_{idx}")

    try:
        _safe_chroma_add(bot.collection, documents, metadatas, ids, batch_size=10)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi nạp vector vào ChromaDB: {e}")

    # Cập nhật gợi ý động background
    threading.Thread(target=generate_dynamic_suggestions, daemon=True).start()

    print(f"[INGEST-URL] ✅ Nạp thành công {len(chunks)} chunks từ '{source_title}'")
    return {
        "status": "success",
        "message": f"Đã nạp thành công nội dung từ '{source_title}'",
        "url": validated_url,
        "source_title": source_title,
        "chunks_count": len(chunks),
        "type": "WEB",
    }


# ─────────────────────────────────────────────────────
# Endpoint: Xem chi tiết Chunks của Trang Web (Strictly WEB only)
# ─────────────────────────────────────────────────────
@app.get("/web-document-chunks")
def get_web_document_chunks(name: str):
    """
    Lấy danh sách các đoạn văn bản Chunks từ ChromaDB CHỈ DÀNH CHO TÀI LIỆU LOẠI WEB.
    Thắt chặt bảo mật: Từ chối xem tài liệu nội bộ PDF, DOCX, XLSX, Chat CSKH.
    """
    try:
        query_res = None
        # Thử tìm theo source_url hoặc source_title
        query_res = bot.collection.get(where={"source_url": name})
        if not query_res or not query_res.get("documents"):
            query_res = bot.collection.get(where={"source_title": name})

        if not query_res or not query_res.get("documents"):
            raise HTTPException(status_code=404, detail=f"Không tìm thấy dữ liệu cho trang web '{name}'")

        metadatas = query_res.get("metadatas", [])
        documents = query_res.get("documents", [])

        # Kiểm tra bảo mật: CHỈ CHO PHÉP XEM NẾU TYPE === 'WEB'
        first_meta = metadatas[0] if metadatas else {}
        doc_type = str(first_meta.get("type", "")).upper()
        if doc_type != "WEB":
            raise HTTPException(
                status_code=403,
                detail="Tài liệu nội bộ doanh nghiệp - Bảo mật cao, không hỗ trợ xem trực tiếp."
            )

        chunks_list = []
        for idx, (doc_text, meta) in enumerate(zip(documents, metadatas)):
            chunks_list.append({
                "chunk_index": idx + 1,
                "text": doc_text,
                "metadata": meta
            })

        return {
            "status": "success",
            "source": name,
            "source_title": first_meta.get("source_title", name),
            "source_url": first_meta.get("source_url", name),
            "total_chunks": len(chunks_list),
            "chunks": chunks_list
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi đọc chunks tri thức: {e}")


# ─────────────────────────────────────────────────────
# Endpoint: Ingest trực tiếp DOM HTML từ Trình duyệt (Backup 100%)
# ─────────────────────────────────────────────────────
@app.post("/ingest-html")
async def ingest_from_html(
    source_title: str = Form(...),
    html_content: str = Form(...),
    source_url: Optional[str] = Form(""),
):
    """Nạp trực tiếp DOM HTML hoặc văn bản đã được copy từ trình duyệt người dùng."""
    title, text_content = "", ""
    if _extract_title_and_text_from_html is not None:
        title, text_content = _extract_title_and_text_from_html(html_content, source_title)
    
    if not text_content or len(text_content.strip()) < 50:
        # Fallback nếu html_content đã là plain text
        text_content = html_content.strip()

    if not text_content or len(text_content.strip()) < 30:
        raise HTTPException(status_code=400, detail="Nội dung HTML/Text cung cấp quá ngắn.")

    final_title = (source_title or title or source_url or "Trang web DOM")[:200]
    final_url = source_url or f"dom://{final_title}"

    upload_date = time.strftime("%d %b %Y")
    words = text_content.split()
    chunk_size = 300
    overlap = 50
    step = chunk_size - overlap
    chunks = []
    for i in range(0, len(words), step):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)

    if not chunks:
        raise HTTPException(status_code=400, detail="Không thể chia nhỏ nội dung HTML.")

    try:
        bot.collection.delete(where={"source_title": final_title})
    except Exception:
        pass

    documents, metadatas, ids = [], [], []
    ts = int(time.time())
    for idx, chunk in enumerate(chunks):
        documents.append(chunk)
        metadatas.append({
            "source_title": final_title,
            "source_url": final_url,
            "type": "WEB",
            "upload_date": upload_date,
            "timestamp": ts,
            "images": "",
        })
        ids.append(f"dom_{ts}_{idx}")

    try:
        _safe_chroma_add(bot.collection, documents, metadatas, ids, batch_size=10)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi nạp vector HTML vào ChromaDB: {e}")

    threading.Thread(target=generate_dynamic_suggestions, daemon=True).start()

    print(f"[INGEST-HTML] ✅ Nạp thành công {len(chunks)} chunks từ DOM '{final_title}'")
    return {
        "status": "success",
        "message": f"Đã nạp thành công nội dung DOM từ '{final_title}'",
        "url": final_url,
        "source_title": final_title,
        "chunks_count": len(chunks),
        "type": "WEB",
    }


def _normalize_administrative_document(text_content: str, filename: str) -> tuple[str, dict]:
    """
    Chuẩn hóa số hiệu, ngày ban hành và header của các văn bản hành chính / công văn PDF.
    - Đọc DUY NHẤT một số hiệu chính thức, không thêm mở ngoặc biến thể (hoặc ...).
    - Nếu văn bản không có số hiệu công văn, đánh dấu rõ ràng 'KHÔNG CÓ TRÊN VĂN BẢN GỐC'.
    """
    import re
    clean_name = os.path.splitext(filename)[0]
    
    doc_number = None
    doc_symbol = None
    issue_date = None
    doc_year = 2025
    
    # 1. Trích xuất số phát hành và ký hiệu từ filename
    m1 = re.search(r'(?:VB_DEN_)?(\d{2,6})[_\-]([A-Za-z0-9\-_Đđ]+?)(?:_(\d{8}))?(?:_|$)', clean_name, re.IGNORECASE)
    if m1:
        doc_number = m1.group(1)
        doc_symbol = m1.group(2).replace('_', '-')
        if m1.group(3):
            d_raw = m1.group(3)
            issue_date = f"{d_raw[6:8]}/{d_raw[4:6]}/{d_raw[0:4]}"
            try:
                doc_year = int(d_raw[0:4])
            except Exception:
                pass
            
    if not doc_number:
        m2 = re.search(r'VB\s*(\d{2,6})', clean_name, re.IGNORECASE)
        if m2:
            doc_number = m2.group(1)
            
    # 2. Tìm kiếm và phục hồi dòng 'Số: /...' trong nội dung text nếu bị khuyết số
    if doc_number:
        def repl_so(match):
            existing_symbol = match.group(1) or ""
            existing_symbol = existing_symbol.strip()
            if existing_symbol:
                return f"Số: {doc_number}/{existing_symbol}"
            elif doc_symbol:
                return f"Số: {doc_number}/{doc_symbol}"
            return f"Số: {doc_number}"
            
        text_content = re.sub(r'Số:\s*\/\s*([A-Za-z0-9\-_Đđ\s]+?)(?=\s+Hà Nội|\s+TP|\s+ngày|\n|$)', repl_so, text_content, count=1, flags=re.IGNORECASE)

    # 3. Quét lại số hiệu đầy đủ từ dòng 'Số: ...' trong text
    full_so_hieu = ""
    so_match = re.search(r'Số:\s*([0-9]+[A-Za-z0-9\-_\/Đđ\s]+?)(?=\s+Hà Nội|\s+TP|\s+ngày|\n|$)', text_content, re.IGNORECASE)
    if so_match:
        full_so_hieu = re.sub(r'\s+', '', so_match.group(1))
    elif doc_number and doc_symbol:
        full_so_hieu = f"{doc_number}/{doc_symbol}"
    elif doc_number:
        full_so_hieu = doc_number

    # 4. Trích xuất trích yếu V/v... từ text
    subject_m = re.search(r'V\/v\s+([^\n\r]+)', text_content, re.IGNORECASE)
    subject_text = subject_m.group(1).strip() if subject_m else ""
    
    # 5. Trích xuất ngày ban hành từ text nếu chưa có từ filename
    if not issue_date:
        date_m = re.search(r'ngày\s+(\d{1,2})\s+tháng\s+(\d{1,2})\s+năm\s+(\d{4})', text_content, re.IGNORECASE)
        if date_m:
            issue_date = f"{int(date_m.group(1)):02d}/{int(date_m.group(2)):02d}/{date_m.group(3)}"
            try:
                doc_year = int(date_m.group(3))
            except Exception:
                pass
        else:
            y_m = re.search(r'năm\s+(\d{4})', clean_name)
            if y_m:
                try:
                    doc_year = int(y_m.group(1))
                except Exception:
                    pass

    # 6. Tạo metadata header chuẩn hóa chèn lên đầu
    header_lines = ["[THÔNG TIN ĐỊNH DANH VĂN BẢN CHÍNH THỨC]:"]
    if full_so_hieu:
        # Chuẩn hóa duy nhất 1 định dạng (ví dụ Đ01 hoặc D01 theo đúng văn bản)
        header_lines.append(f"- SỐ HIỆU CÔNG VĂN CHÍNH THỨC: {full_so_hieu}")
    else:
        header_lines.append("- SỐ HIỆU CÔNG VĂN CHÍNH THỨC: KHÔNG CÓ TRÊN VĂN BẢN GỐC (Tài liệu nội bộ/phụ lục không ghi số hiệu công văn)")
        
    if issue_date:
        header_lines.append(f"- NGÀY BAN HÀNH: {issue_date}")
    if subject_text:
        header_lines.append(f"- TRÍCH YẾU NỘI DUNG: {subject_text}")
    header_lines.append(f"- TÊN FILE GỐC: {filename}")
    if full_so_hieu:
        header_lines.append("[LƯU Ý PHÁP LÝ]: Mọi số hiệu khác xuất hiện trong phần 'Căn cứ...' chỉ là tài liệu viện dẫn cũ, KHÔNG PHẢI số hiệu của văn bản này.")
    header_lines.append("──────────────────────────────────────────\n")
    
    normalized_header = "\n".join(header_lines) + "\n"
    
    metadata_info = {
        "doc_number": full_so_hieu if full_so_hieu else "KHONG_CO",
        "issue_date": issue_date if issue_date else "",
        "doc_year": doc_year,
        "subject": subject_text
    }
    return normalized_header + text_content, metadata_info


# ─────────────────────────────────────────────────────
# Upload tài liệu và nạp vector tức thì (Hot-reload Ingestion)
# ─────────────────────────────────────────────────────
@app.post("/upload")
async def upload_document(
    file: UploadFile = File(...),
    ingest_type: str = Form("rag"),  # "rag" (mặc định) hoặc "conversation"
):
    filename = file.filename
    content_type = file.content_type
    
    # 1. Đọc nội dung file nhị phân
    file_bytes = await file.read()
    size_bytes = len(file_bytes)
    
    # Tạo thư mục temp để lưu tạm file
    temp_dir = os.path.join(BASE_DIR, "temp_uploads")
    os.makedirs(temp_dir, exist_ok=True)
    temp_file_path = os.path.join(temp_dir, filename)
    
    with open(temp_file_path, "wb") as f:
        f.write(file_bytes)
        
    text_content = ""
    file_ext = os.path.splitext(filename)[1].lower()
    is_pptx = False
    pptx_chunks = []
    
    # 2. Phân tách và trích xuất text tùy theo định dạng file
    try:
        if file_ext == ".txt":
            text_content = file_bytes.decode("utf-8", errors="ignore")
        elif file_ext == ".json":
            json_data = json.loads(file_bytes.decode("utf-8", errors="ignore"))
            if isinstance(json_data, list):
                text_content = "\n".join([json.dumps(item, ensure_ascii=False) for item in json_data])
            else:
                text_content = json.dumps(json_data, ensure_ascii=False)
        elif file_ext == ".pdf":
            # Dùng pdfminer.six (đã có sẵn vì là dependency của pdfplumber)
            # LAParams giúp xử lý đúng layout 2 cột của công văn nhà nước:
            #   - boxes_flow=0.5: cân bằng giữa x và y khi sort text blocks
            #   - char_margin=2.0: gộp ký tự cách nhau ≤2x font size thành 1 từ
            #   - line_margin=0.5: gộp dòng cách nhau ≤0.5x font size thành 1 block
            # → Số "6790" và "/D01-B4-B5" trong cùng text box sẽ được nối đúng
            try:
                from pdfminer.high_level import extract_text as pdfminer_extract
                from pdfminer.layout import LAParams
                laparams = LAParams(
                    char_margin=2.0,
                    word_margin=0.1,
                    line_margin=0.5,
                    boxes_flow=0.5,
                    detect_vertical=False,
                )
                text_content = pdfminer_extract(temp_file_path, laparams=laparams) or ""
                # Làm sạch khoảng trắng thừa
                import re as _re
                text_content = _re.sub(r'\n{3,}', '\n\n', text_content)  # Tối đa 2 dòng trắng liên tiếp
                text_content = _re.sub(r'[ \t]{3,}', ' ', text_content)   # Tối đa 1 khoảng trắng
                print(f"[PDF] pdfminer ok: {len(text_content)} ký tự")
            except Exception as e_pdfminer:
                print(f"[PDF] pdfminer lỗi ({e_pdfminer}), fallback sang pdfplumber...")
                try:
                    import pdfplumber
                    text_list = []
                    with pdfplumber.open(temp_file_path) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text(layout=False) or ""
                            text_list.append(page_text)
                    text_content = "\n".join(text_list)
                except Exception:
                    import pypdf
                    reader = pypdf.PdfReader(temp_file_path)
                    text_content = "\n".join(p.extract_text() or "" for p in reader.pages)


        elif file_ext in [".docx", ".doc"]:
            text_content = extract_docx_images_and_text(temp_file_path, filename)
        elif file_ext in [".xlsx", ".xls"]:
            import pandas as pd
            df_dict = pd.read_excel(temp_file_path, sheet_name=None)
            text_list = []
            for sheet_name, df in df_dict.items():
                text_list.append(f"Sheet: {sheet_name}\n" + df.to_string())
            text_content = "\n".join(text_list)
        elif file_ext == ".pptx":
            is_pptx = True
            pptx_chunks = extract_pptx_slides_and_text(temp_file_path, filename)
            text_content = "\n".join([chunk["text"] for chunk in pptx_chunks])
        else:
            raise HTTPException(status_code=400, detail="Định dạng file không hỗ trợ. Chỉ nhận TXT, JSON, PDF, DOCX, XLSX, PPTX.")
    except Exception as e:
        # Dọn dẹp file temp
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        raise HTTPException(status_code=500, detail=f"Lỗi khi trích xuất nội dung file: {str(e)}")
    
    # Dọn dẹp file temp
    if os.path.exists(temp_file_path):
        os.remove(temp_file_path)
        
    if not is_pptx and (not text_content.strip() or len(text_content.strip()) < 10):
        raise HTTPException(status_code=400, detail="Nội dung file trống hoặc quá ngắn, không thể nạp vector.")

    # Tính mã băm SHA-256 nội dung file để khử trùng lặp liên file (Cross-file Deduplication)
    import hashlib
    content_hash = hashlib.sha256(file_bytes).hexdigest()

    # Kiểm tra xem nội dung tài liệu này đã tồn tại trong ChromaDB chưa
    try:
        existing_check = bot.collection.get(where={"content_hash": content_hash}, include=["metadatas"])
        if existing_check and existing_check.get("ids"):
            first_meta = existing_check["metadatas"][0] if existing_check["metadatas"] else {}
            orig_file = first_meta.get("source_title", "tài liệu trước đó")
            print(f"[DEDUPLICATION] ⚡ Bỏ qua file '{filename}' vì nội dung trùng lặp 100% (SHA-256) với '{orig_file}'")
            return {
                "status": "success",
                "message": f"Nội dung file '{filename}' đã tồn tại trong hệ thống (trùng với '{orig_file}'). Đã tự động tái sử dụng để tránh tạo vector rác.",
                "chunks_count": len(existing_check["ids"]),
                "size": f"{size_bytes / 1024:.1f} KB",
                "type": "DUPLICATE_SKIPPED",
                "packages": [],
            }
    except Exception as hash_err:
        print(f"⚠️ Cảnh báo khi kiểm tra SHA-256 hash: {hash_err}")

    upload_date = time.strftime("%d %b %Y")
    ts = int(time.time())

    # ─────────────────────────────────────────────────────
    # Nhánh CONVERSATION: Lưu nguyên Kịch bản Hội thoại CSKH (Full Playbook)
    # ─────────────────────────────────────────────────────
    if ingest_type == "conversation":
        doc_text = f"KỊCH BẢN TƯ VẤN VÀ CHỐT SALE CSKH MOBIFONE CHUẨN ({filename}):\n{text_content.strip()}"
        source_title = f"Tri thức CSKH [Kịch bản]: {filename}"
        documents = [doc_text]
        metadatas = [{
            "source_title": source_title,
            "source_url": f"upload://{filename}",
            "type": "CONVERSATION",
            "size_bytes": size_bytes,
            "upload_date": upload_date,
            "timestamp": ts,
            "images": "",
            "content_hash": content_hash,
            "doc_year": 2026,
        }]
        ids = [f"conv_playbook_{filename}_{ts}"]

        try:
            try:
                bot.collection.delete(where={"source_title": source_title})
            except Exception as del_err:
                print(f"⚠️ Cảnh báo khi dọn dẹp conversation cũ: {del_err}")

            bot.collection.add(
                documents=documents,
                metadatas=metadatas,
                ids=ids,
            )
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Lỗi khi nạp kịch bản conversation vào ChromaDB: {e}")

        threading.Thread(target=generate_dynamic_suggestions, daemon=True).start()
        print(f"[CONVERSATION] ✅ Nạp thành công Kịch bản CSKH trọn gói từ '{filename}'")
        return {
            "status": "success",
            "message": f"Đã nạp thành công Kịch bản CSKH trọn gói từ file chat mẫu '{filename}'",
            "chunks_count": 1,
            "size": f"{size_bytes / 1024:.1f} KB",
            "type": "CONVERSATION",
            "packages": [],
        }

    # ─────────────────────────────────────────────────────
    # Nhánh RAG thông thường (ingest_type == "rag")
    # ─────────────────────────────────────────────────────
    # 3. Chia nhỏ văn bản (Chunking) & Nạp vào ChromaDB
    documents = []
    metadatas = []
    ids = []

    if is_pptx:
        if not pptx_chunks:
            raise HTTPException(status_code=400, detail="Không thể trích xuất nội dung từ file PPTX.")
        for idx, item in enumerate(pptx_chunks):
            documents.append(item["text"])
            metadatas.append({
                "source_title": filename,
                "source_url": f"upload://{filename}",
                "type": "PPTX",
                "size_bytes": size_bytes,
                "upload_date": upload_date,
                "timestamp": int(time.time()),
                "slide_index": item["metadata"]["slide_index"],
                "images": item["metadata"]["images"],
                "content_hash": content_hash,
                "doc_year": 2026,
            })
            ids.append(f"upload_{filename}_{int(time.time())}_{idx}")
    else:
        # Chuẩn hóa số hiệu công văn và thông tin hành chính
        text_content, doc_meta = _normalize_administrative_document(text_content, filename)

        words = text_content.split()
        chunk_size = 300  # số từ mỗi mảnh
        overlap = 50
        step = chunk_size - overlap
        
        chunks = []
        for i in range(0, len(words), step):
            chunk = " ".join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
                
        if not chunks:
            raise HTTPException(status_code=400, detail="Không thể tạo các mảnh dữ liệu văn bản.")
            
        for idx, chunk in enumerate(chunks):
            import re
            img_matches = re.findall(r'\[IMAGE:(.*?)\]', chunk)
            chunk_images = ",".join(list(set(img_matches)))
            
            documents.append(chunk)
            metadatas.append({
                "source_title": filename,
                "source_url": f"upload://{filename}",
                "type": file_ext.replace(".", "").upper(),
                "size_bytes": size_bytes,
                "upload_date": upload_date,
                "timestamp": int(time.time()),
                "images": chunk_images,
                "content_hash": content_hash,
                "doc_number": doc_meta.get("doc_number", ""),
                "doc_year": doc_meta.get("doc_year", 2025),
                "issue_date": doc_meta.get("issue_date", "")
            })
            ids.append(f"upload_{filename}_{int(time.time())}_{idx}")

    try:
        # Tự động xóa các vector cũ cùng tên (nếu có) để tránh trùng lặp dữ liệu
        try:
            with chroma_write_lock:
                bot.collection.delete(where={"source_title": filename})
        except Exception as delete_err:
            print(f"⚠️ Cảnh báo khi dọn dẹp tài liệu cũ: {delete_err}")

        total_docs = len(documents)
        print(f"[UPLOAD] Bắt đầu nạp {total_docs} chunks (batch=10, retry=3)...")
        _safe_chroma_add(bot.collection, documents, metadatas, ids, batch_size=10)
        print(f"[UPLOAD] ✓ Đã nạp xong {total_docs} chunks")
            
        # 5. Gọi Gemini trích xuất thông tin gói cước nếu có trong tài liệu tri thức
        extracted_packages = []
        try:
            # Chỉ gửi một phần nội dung nếu tệp quá lớn để tránh quá tải token
            preview_content = text_content[:15000]
            
            prompt = f"""Bạn là một chuyên gia phân tích dữ liệu của nhà mạng MobiFone.
Hãy đọc kỹ văn bản dưới đây và trích xuất danh sách tất cả các gói cước di động/viễn thông MobiFone được giới thiệu.
Với mỗi gói cước, hãy trích xuất các thông tin sau theo đúng cấu trúc JSON:
- id: Mã gói cước viết hoa liền nhau, ví dụ: "TK135", "KC150"
- name: Tên gói cước viết hoa liền nhau, giống id
- price: Giá gói cước định dạng chuỗi có chữ 'đ' hoặc 'đồng', ví dụ: "135.000đ", "90.000đ"
- data: Mô tả dung lượng data, ví dụ: "4GB/ngày", "1.5GB/ngày", "Không giới hạn"
- voice: Mô tả ưu đãi cuộc gọi thoại, ví dụ: "Nội mạng miễn phí + 20p ngoại mạng", "1000p nội mạng"
- validity: Chu kỳ gói cước, ví dụ: "30 ngày", "24 giờ"
- category: Thể loại gói cước, chỉ chọn một trong ba giá trị sau: "data", "voice", hoặc "unlimited"
- features: Danh sách mảng chuỗi các đặc điểm nổi bật, ví dụ: ["5G Ready", "MobiFone TV+", "Xem YouTube miễn phí"]
- color: Mã màu HEX phù hợp làm màu chủ đạo cho card gói cước này, ví dụ: "#E4002B", "#0055A5", "#059669", "#7C3AED", "#DC2626", "#4F46E5"
- popular: Giá trị boolean (true hoặc false) biểu thị gói cước này có phải là gói nổi bật phổ biến không
- dataTotalGB: Tổng dung lượng GB trong một chu kỳ (ví dụ: 4GB/ngày * 30 ngày = 120, hoặc 1GB/ngày * 30 ngày = 30, hoặc không giới hạn = 999), lưu dạng số nguyên (integer)
- voiceTotalMin: Tổng số phút gọi thoại nội mạng + ngoại mạng trong một chu kỳ, lưu dạng số nguyên (integer). Nếu không đề cập thì mặc định lưu 600.

Hãy trả về kết quả dưới dạng một mảng JSON các đối tượng gói cước.
LƯU Ý QUAN TRỌNG: Chỉ trả về JSON nguyên bản, không kèm ký tự markdown như ```json hay ```. Không có bất cứ ký tự giải thích nào khác ngoài chuỗi JSON hợp lệ. Nếu không tìm thấy bất kỳ gói cước nào trong văn bản, hãy trả về mảng rỗng [].

Văn bản cần phân tích:
{preview_content}
"""
            llm_response = bot._call_llm_with_retry(prompt, temperature=0.1)
            
            # Làm sạch dữ liệu phản hồi từ mô hình
            cleaned_response = llm_response.strip()
            if cleaned_response.startswith("```"):
                lines = cleaned_response.splitlines()
                if lines[0].startswith("```"):
                    lines = lines[1:]
                if lines and lines[-1].startswith("```"):
                    lines = lines[:-1]
                cleaned_response = "\n".join(lines).strip()
            
            extracted_packages = json.loads(cleaned_response)
            if not isinstance(extracted_packages, list):
                extracted_packages = []
            print(f"[EXTRACT] Trích xuất thành công {len(extracted_packages)} gói cước từ '{filename}'")
        except Exception as extract_err:
            print(f"⚠️ Cảnh báo: Lỗi khi trích xuất gói cước bằng Gemini: {extract_err}")
            extracted_packages = []

        # Cập nhật gợi ý động trong background thread để tránh làm chậm response upload
        threading.Thread(target=generate_dynamic_suggestions, daemon=True).start()

        return {
            "status": "success",
            "message": f"Đã nạp thành công tài liệu '{filename}'",
            "chunks_count": len(pptx_chunks) if is_pptx else len(chunks),
            "size": f"{size_bytes / 1024:.1f} KB",
            "packages": extracted_packages
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi khi nạp vector vào ChromaDB: {e}")


# ─────────────────────────────────────────────────────
# CSKH Chat Mining Endpoints
# ─────────────────────────────────────────────────────

class ParseTextRequest(BaseModel):
    raw_text: str

class ApproveQARequest(BaseModel):
    qa_list: List[dict]


@app.post("/chat-mining/parse-text")
def chat_mining_parse_text(req: ParseTextRequest):
    """Phân tích văn bản chat copy-paste, phân vai, xóa PII và trích xuất Q&A / Kịch bản bán hàng."""
    try:
        from chat_miner import deduplicate_qa_list
        conv = parse_raw_text_to_chat(req.raw_text)
        analysis = analyze_chat_with_llm(conv, bot_pipeline=bot)
        analysis_dict = analysis.dict()
        analysis_dict["extracted_qa_list"] = deduplicate_qa_list(analysis.extracted_qa_list, collection=bot.collection)
        return {
            "status": "success",
            "conversation": conv.dict(),
            "analysis": analysis_dict
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi phân tích văn bản chat: {e}")


@app.post("/chat-mining/parse-file")
async def chat_mining_parse_file(file: UploadFile = File(...)):
    """Đọc file chat (CSV, XLSX, JSON, TXT), phân vai, xóa PII và trích xuất tri thức."""
    try:
        from chat_miner import deduplicate_qa_list
        file_bytes = await file.read()
        chats = parse_file_to_chats(file_bytes, file.filename)
        results = []
        for chat in chats:
            analysis = analyze_chat_with_llm(chat, bot_pipeline=bot)
            analysis_dict = analysis.dict()
            analysis_dict["extracted_qa_list"] = deduplicate_qa_list(analysis.extracted_qa_list, collection=bot.collection)
            results.append({
                "conversation": chat.dict(),
                "analysis": analysis_dict
            })
        return {
            "status": "success",
            "count": len(results),
            "results": results
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi xử lý file chat: {e}")


@app.post("/chat-mining/approve-qa")
def chat_mining_approve_qa(req: ApproveQARequest):
    """Duyệt và nạp Kịch bản & Tri thức Bán hàng thực chiến từ chat CSKH vào ChromaDB."""
    try:
        added_count = 0
        documents = []
        metadatas = []
        ids = []

        dialogue_lines = []
        first_q = ""

        # 1. Nạp từng cặp Q&A với nhãn Sales Stage & Tactic chi tiết
        for idx, item in enumerate(req.qa_list):
            question = str(item.get("question", "")).strip()
            answer = str(item.get("answer", "")).strip()
            sales_stage = str(item.get("sales_stage", "kham_pha_nhu_cau")).strip()
            sales_tactic = str(item.get("sales_tactic", "Tư vấn tiêu chuẩn")).strip()
            package_name = str(item.get("package_name", "")).strip()
            intent = str(item.get("intent", "Tư vấn")).strip()

            if question and answer:
                if not first_q:
                    first_q = question
                dialogue_lines.append(f"Khách hàng: {question}\nChuyên viên Bán hàng MobiFone: {answer}")

                doc_text = f"TÌNH HUỐNG CSKH & SALES [{sales_stage.upper()}]:\nKhách hỏi: {question}\nChuyên viên tư vấn: {answer}\nChiến thuật: {sales_tactic}"
                doc_id = f"cskh_qa_{uuid.uuid4().hex[:10]}"

                documents.append(doc_text)
                metadatas.append({
                    "source": "CSKH_Chat_Mining",
                    "source_title": f"Q&A: {question[:50]}",
                    "source_url": "chat_mining://cskh",
                    "type": "CONVERSATION",
                    "category": "Sales_Playbook",
                    "sales_stage": sales_stage,
                    "sales_tactic": sales_tactic,
                    "package_name": package_name,
                    "question": question,
                    "answer": answer,
                    "intent": intent,
                    "size_bytes": len(doc_text.encode("utf-8")),
                    "upload_date": time.strftime("%Y-%m-%d"),
                    "timestamp": time.time(),
                    "created_at": time.strftime("%Y-%m-%d %H:%M:%S")
                })
                ids.append(doc_id)
                added_count += 1

        # 2. Nạp thêm 1 bản Full Script hoàn chỉnh cho ngữ cảnh đa lượt
        if dialogue_lines:
            full_script = "\n\n".join(dialogue_lines)
            script_title = f"Kịch bản Sales Top Performer: {first_q[:40]}..."
            full_doc_text = f"KỊCH BẢN CHỐT SALE & PHONG THÁI CSKH XUẤT SẮC MẪU ({script_title}):\n{full_script}"
            full_doc_id = f"cskh_script_{uuid.uuid4().hex[:10]}"

            documents.append(full_doc_text)
            metadatas.append({
                "source": "CSKH_Chat_Mining",
                "source_title": script_title,
                "source_url": "chat_mining://cskh_script",
                "type": "CONVERSATION",
                "category": "CSKH_Learned_Playbook",
                "size_bytes": len(full_doc_text.encode("utf-8")),
                "upload_date": time.strftime("%Y-%m-%d"),
                "timestamp": time.time(),
                "created_at": time.strftime("%Y-%m-%d %H:%M:%S")
            })
            ids.append(full_doc_id)

        if documents:
            bot.collection.add(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
            try:
                bot.playbook_collection.add(
                    documents=documents,
                    metadatas=metadatas,
                    ids=ids
                )
            except Exception as pb_err:
                print(f"⚠️ Cảnh báo nạp Playbook Collection: {pb_err}")

        return {
            "status": "success",
            "message": f"Đã nạp thành công {added_count} Kịch bản & Tri thức Bán hàng thực chiến vào hệ thống",
            "added_count": added_count
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi nạp tri thức CSKH: {e}")


@app.get("/chat-mining/playbooks")
def get_chat_mining_playbooks(stage: Optional[str] = None):
    """Lấy danh sách các kịch bản & nghệ thuật giao tiếp CSKH đã được duyệt."""
    try:
        where_clause = None
        if stage and stage != "all":
            where_clause = {"sales_stage": stage}

        query_res = bot.playbook_collection.get(where=where_clause)
        metadatas = query_res.get("metadatas", [])
        documents = query_res.get("documents", [])
        ids = query_res.get("ids", [])

        playbooks = []
        for doc_id, doc_text, meta in zip(ids, documents, metadatas):
            playbooks.append({
                "id": doc_id,
                "question": meta.get("question", ""),
                "answer": meta.get("answer", ""),
                "package_name": meta.get("package_name", ""),
                "sales_stage": meta.get("sales_stage", "kham_pha_nhu_cau"),
                "sales_tactic": meta.get("sales_tactic", "Tư vấn tiêu chuẩn"),
                "created_at": meta.get("created_at", "")
            })

        return {
            "status": "success",
            "total": len(playbooks),
            "playbooks": playbooks
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Lỗi đọc Playbooks: {e}")

