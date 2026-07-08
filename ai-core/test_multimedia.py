import os
import sys
import shutil

if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')

# Tạo file ảnh giả lập bằng Pillow
from PIL import Image
img = Image.new('RGB', (100, 100), color='red')
img.save("dummy.png")

print("✅ Đã tạo file dummy.png")

# 1. Tạo file docx giả lập chứa ảnh
import docx
doc = docx.Document()
doc.add_paragraph("Đoạn văn bản có chứa ảnh hướng dẫn đăng ký eSIM.")
doc.add_picture("dummy.png")
doc.add_paragraph("Vui lòng làm theo hình ảnh trên.")
doc.save("dummy_test.docx")
print("✅ Đã tạo file dummy_test.docx")

# 2. Tạo file pptx giả lập chứa ảnh
from pptx import Presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
slide.shapes.add_picture("dummy.png", 0, 0)
txBox = slide.shapes.add_textbox(0, 0, 100, 100)
txBox.text_frame.text = "Slide hướng dẫn cài đặt gói cước 4G MobiFone."
prs.save("dummy_test.pptx")
print("✅ Đã tạo file dummy_test.pptx")

# 3. Test trích xuất từ docx
from api_server import extract_docx_images_and_text, extract_pptx_slides_and_text
print("\n--- Test Trích xuất Docx ---")
docx_text = extract_docx_images_and_text("dummy_test.docx", "dummy_test.docx")
print("Text trích xuất được:")
print(docx_text)
assert "[IMAGE:" in docx_text, "Lỗi: Không tìm thấy placeholder ảnh trong Docx!"
print("✅ Test Docx thành công!")

# 4. Test trích xuất từ pptx
print("\n--- Test Trích xuất Pptx ---")
pptx_chunks = extract_pptx_slides_and_text("dummy_test.pptx", "dummy_test.pptx")
print("Các slide trích xuất được:")
for idx, chunk in enumerate(pptx_chunks):
    print(f"Slide {idx + 1}: {chunk}")
    assert "[IMAGE:" in chunk["text"], "Lỗi: Không tìm thấy placeholder ảnh trong slide PPTX!"
print("✅ Test Pptx thành công!")

# 5. Dọn dẹp các file dummy
os.remove("dummy.png")
os.remove("dummy_test.docx")
os.remove("dummy_test.pptx")
print("\n✅ Đã dọn dẹp các file tạm.")
print("=" * 60)
print("TẤT CẢ MULTIMEDIA SMOKE TESTS ĐÃ VƯỢT QUA!")
print("=" * 60)
